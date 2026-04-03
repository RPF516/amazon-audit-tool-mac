# app/services/ppt_service.py

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


import os
import sys
from pptx import Presentation

def get_base_path():
    if getattr(sys, 'frozen', False):
        return sys._MEIPASS  # PyInstaller temp folder
    return os.path.dirname(os.path.abspath(__file__))


def generate_ppt(data, target_acos, search_data, brand_name=""):
    BASE_PATH = get_base_path()

    template_path = os.path.join(BASE_PATH, "template.pptx")
    output_path = os.path.join(os.getcwd(), "output.pptx")

    prs = Presentation(template_path)

    # Slide 1
    replace_slide1(prs.slides[0].shapes, brand_name)
    
    # Slide 3
    replace_slide3(prs.slides[2].shapes, data, target_acos)

    # Slide 4
    replace_slide4(prs.slides[3].shapes, data, target_acos)

    # Slide 5
    replace_slide5(prs.slides[4], data)
    
    # Slide 6
    replace_slide6(prs.slides[5], search_data)
    
    # Slide 7
    replace_slide7(prs.slides[6], search_data)
    
    # Slide 8
    replace_slide8(prs.slides[7], search_data)
    
    # Slide 9
    replace_slide9(prs.slides[8], search_data)
    
    # Slide 10
    replace_slide10(prs.slides[9], search_data)

    prs.save(output_path)

    return output_path
    
from pptx.util import Pt

def replace_slide1(shapes, brand_name):
    for shape in shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            # Check full paragraph text first
            full_text = paragraph.text
            if "{{brand_name}}" not in full_text:
                continue

            # If placeholder is split across runs, merge into first run
            if len(paragraph.runs) > 1:
                first_run = paragraph.runs[0]
                first_run.text = full_text.replace("{{brand_name}}", brand_name)
                for run in paragraph.runs[1:]:
                    run.text = ""
            elif len(paragraph.runs) == 1:
                paragraph.runs[0].text = full_text.replace("{{brand_name}}", brand_name)
            else:
                paragraph.text = full_text.replace("{{brand_name}}", brand_name)

def replace_slide3(shapes, data, target_acos):
    replacements = {
        "{{profitable_keywords}}": str(data["TOTALS"]["profitable_keywords"]),
        "{{unprofitable_keywords}}": str(data["TOTALS"]["unprofitable_keywords"]),
        "{{clicks_no_sales}}": str(data["TOTALS"]["clicks_no_sales"]),
        "{{zero_clicks}}": str(data["TOTALS"]["zero_clicks"]),
        "{{target_acos}}": str(int(target_acos)),
    }

    for shape in shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            original_text = paragraph.text
            new_text = original_text

            for key, value in replacements.items():
                new_text = new_text.replace(key, value)

            # Only update if something changed
            if new_text != original_text:
                paragraph.text = new_text

                # 🔥 Resize ONLY numbers (not ACOS)
                for run in paragraph.runs:
                    if any(val in run.text for val in [
                        str(data["TOTALS"]["profitable_keywords"]),
                        str(data["TOTALS"]["unprofitable_keywords"]),
                        str(data["TOTALS"]["clicks_no_sales"]),
                        str(data["TOTALS"]["zero_clicks"]),
                    ]):
                        run.font.size = Pt(60)
        
def replace_slide4(shapes, data, target_acos):
    replacements = {
        "{{profitable_sales}}": str(round(data["TOTALS"]["profitable_sales"], 2)),
        "{{unprofitable_sales}}": str(round(data["TOTALS"]["unprofitable_sales"], 2)),
        "{{wasted_spend}}": str(round(data["TOTALS"]["wasted_spend"], 2)),
        "{{target_acos}}": str(int(target_acos)),
    }

    for shape in shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            original_text = paragraph.text
            new_text = original_text

            for key, value in replacements.items():
                new_text = new_text.replace(key, value)

            if new_text != original_text:
                paragraph.text = new_text

                for run in paragraph.runs:
                    if any(val in run.text for val in [
                        str(round(data["TOTALS"]["profitable_sales"], 2)),
                        str(round(data["TOTALS"]["unprofitable_sales"], 2)),
                        str(round(data["TOTALS"]["wasted_spend"], 2)),
                    ]):
                        run.font.size = Pt(54)

def apply_table_font(cell, font_size):
    tf = cell.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = font_size


def fill_match_type_table(table, data):
    col_keys = [
        "Broad",
        "Phrase",
        "Exact",
        "ASIN Targeting",
        "Category Targeting"
    ]

    font_size = Pt(28)  # 🔥 adjust (16–20 works best)

    # Row 1 → Number of Keywords
    for i, key in enumerate(col_keys):
        value = data.get(key, {}).get("total_keywords", 0)
        cell = table.cell(1, i + 1)
        cell.text = str(value)
        apply_table_font(cell, font_size)

    # Row 2 → ACOS
    for i, key in enumerate(col_keys):
        value = data.get(key, {}).get("acos", 0)
        cell = table.cell(2, i + 1)
        cell.text = f"{round(value, 0)}%"
        apply_table_font(cell, font_size)

    # Row 3 → Sales
    for i, key in enumerate(col_keys):
        value = data.get(key, {}).get("sales", 0)
        cell = table.cell(3, i + 1)
        cell.text = f"${round(value, 2)}"
        apply_table_font(cell, font_size)

    # Row 4 → Wasted Spend
    for i, key in enumerate(col_keys):
        value = data.get(key, {}).get("wasted_spend", 0)
        cell = table.cell(4, i + 1)
        cell.text = f"${round(value, 2)}"
        apply_table_font(cell, font_size)
   
def replace_slide5(slide, data):
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            fill_match_type_table(table, data)
            return

    print("❌ No table found in slide 5")
                        
def replace_text_in_shapes(shapes, replacements):
    for shape in shapes:
        if not shape.has_text_frame:
            continue

        # ✅ First try full text replace (works in most cases)
        full_text = shape.text
        new_text = full_text

        for key, value in replacements.items():
            new_text = new_text.replace(key, str(value))

        if new_text != full_text:
            shape.text = new_text
            continue  # done, move to next shape

        # ✅ Fallback → handle split runs
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for key, value in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, str(value))             

def style_replaced_text(shapes, replacements, font_size=Pt(32), color_map=None):
    for shape in shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            original_text = paragraph.text
            new_text = original_text

            # Replace text
            for key, value in replacements.items():
                new_text = new_text.replace(key, str(value))

            if new_text != original_text:
                paragraph.text = new_text

                # Apply styling ONLY to replaced values
                for run in paragraph.runs:
                    for key, value in replacements.items():
                        if str(value) in run.text:
                            run.font.size = font_size

                            if color_map and key in color_map:
                                run.font.color.rgb = color_map[key]  
  

def update_campaign_pie_chart(slide, search_data):
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart

            chart_data = CategoryChartData()

            chart_data.categories = [
                "Below Target",
                "Above Target",
                "Non-Converting"
            ]

            chart_data.add_series(
                "Campaigns",
                [
                    search_data["converting_below_target_count"],
                    search_data["converting_above_target_count"],
                    search_data["non_converting_campaigns_count"]
                ]
            )

            chart.replace_data(chart_data)


def replace_slide6(slide, search_data):
    replacements = {
        "{{converting_below_target_count}}": search_data["converting_below_target_count"],
        "{{converting_above_target_count}}": search_data["converting_above_target_count"],
        "{{non_converting_campaigns_count}}": search_data["non_converting_campaigns_count"],
    }

    color_map = {
        "{{converting_below_target_count}}": RGBColor(255, 255, 255),   # green
        "{{converting_above_target_count}}": RGBColor(255, 255, 255),  # yellow
        "{{non_converting_campaigns_count}}": RGBColor(255, 255, 255),   # red
    }

    style_replaced_text(
        slide.shapes,
        replacements,
        font_size=Pt(24),
        color_map=color_map
    )

    update_campaign_pie_chart(slide, search_data)
            
   

def get_table_font_size(rows):
    # Find longest campaign name
    max_length = 0

    for row in rows:
        campaign = str(row.get("Campaign Name", ""))
        max_length = max(max_length, len(campaign))

    # Decide font size based on longest text
    if max_length < 100:
        return Pt(18)
    elif max_length < 40:
        return Pt(20)
    elif max_length < 70:
        return Pt(22)
    else:
        return Pt(16)


def fill_keyword_table(table, rows):
    # 🔥 Step 1: get ONE font size for entire table
    font_size = get_table_font_size(rows)

    for i, row_data in enumerate(rows):
        row = table.rows[i + 1]

        values = [
            row_data.get("Campaign Name", ""),
            row_data.get("Ad Group Name", ""),
            row_data.get("Customer Search Term", ""),
            f"${round(row_data.get('Spend', 0), 2)}",
            f"${round(row_data.get('7 Day Total Sales', 0), 2)}",
            f"{round(row_data.get('acos', 0), 2)}%",
        ]

        for j, value in enumerate(values):
            cell = row.cells[j]
            cell.text = str(value)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE

            # 🔥 Apply SAME font size everywhere
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size
    
def replace_slide7(slide, search_data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "LOWEST_ACOS_TABLE" in shape.text:

                # find table in same slide
                for shp in slide.shapes:
                    if shp.has_table:
                        table = shp.table

                        fill_keyword_table(
                            table,
                            search_data["lowest_acos_keywords"]
                        )
                        return
                     
def replace_slide8(slide, search_data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "HIGHEST_ACOS_TABLE" in shape.text:

                # find table
                for shp in slide.shapes:
                    if shp.has_table:
                        table = shp.table

                        fill_keyword_table(
                            table,
                            search_data["highest_acos_keywords"]
                        )
                        return
                        
                        
def fill_clicks_only_table(table, rows):
    font_size = get_table_font_size(rows)

    for i, row_data in enumerate(rows):
        row = table.rows[i + 1]

        values = [
            row_data.get("Campaign Name", ""),
            row_data.get("Ad Group Name", ""),
            row_data.get("Customer Search Term", ""),
            str(row_data.get("Clicks", 0)),
        ]

        for j, value in enumerate(values):
            cell = row.cells[j]
            cell.text = str(value)

            tf = cell.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE

            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.size = font_size
        
def replace_slide9(slide, search_data):
    for shape in slide.shapes:
        if shape.has_text_frame:
            if "HIGHEST_CLICK_KW" in shape.text:

                for shp in slide.shapes:
                    if shp.has_table:
                        table = shp.table

                        fill_clicks_only_table(
                            table,
                            search_data["highest_click_no_sales_keywords"]
                        )
                        return
                        
                        

def update_duplicate_pie_chart(slide, search_data):
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart

            chart_data = CategoryChartData()

            chart_data.categories = [
                "Total Duplicates",
                "Unique Duplicates",
                "Exact+Phrase"
            ]

            chart_data.add_series(
                "Keywords",
                [
                    search_data["total_duplicate_keywords"],
                    search_data["unique_duplicate_keywords"],
                    search_data["exact_phrase_duplicate_keywords"]
                ]
            )

            chart.replace_data(chart_data)

def replace_slide10(slide, search_data):
    replacements = {
        "{{total_duplicate_keywords}}": search_data["total_duplicate_keywords"],
        "{{unique_duplicate_keywords}}": search_data["unique_duplicate_keywords"],
        "{{exact_phrase_duplicate_keywords}}": search_data["exact_phrase_duplicate_keywords"],
    }

    color_map = {
        "{{total_duplicate_keywords}}": RGBColor(255, 255, 255),   # red
        "{{unique_duplicate_keywords}}": RGBColor(255, 255, 255), # blue
        "{{exact_phrase_duplicate_keywords}}": RGBColor(255, 255, 255), # purple
    }

    style_replaced_text(
        slide.shapes,
        replacements,
        font_size=Pt(24),
        color_map=color_map
    )

    update_duplicate_pie_chart(slide, search_data)
            